from pptx import Presentation
from pptx.util import Inches, Pt

def create_ppt(title, slides, output_file, max_words_per_slide=130):
    prs = Presentation()
    # الشريحة الأولى - العنوان الرئيسي
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = "Generated Automatically"

    for item in slides:
        slide_title = item.get("title", "")
        subtitle = item.get("subtitle", "")
        content = item.get("content", [])

        current_content = []
        current_word_count = 0

        def add_slide(content_points, show_subtitle):
            slide_layout = prs.slide_layouts[1]  # Title and Content
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = slide_title

            if show_subtitle and subtitle:
                left, top, width, height = Inches(0.5), Inches(1), Inches(9), Inches(0.5)
                textbox = slide.shapes.add_textbox(left, top, width, height)
                tf = textbox.text_frame
                p = tf.add_paragraph()
                p.text = subtitle
                p.font.size = Pt(16)
                p.font.italic = True

            content_placeholder = slide.placeholders[1]
            content_placeholder.left = Inches(1)
            content_placeholder.top = Inches(1)
            content_placeholder.width = Inches(8)
            content_placeholder.height = Inches(5)

            tf = content_placeholder.text_frame
            tf.clear()
            tf.word_wrap = True

            for point in content_points:
                p = tf.add_paragraph()
                p.text = point
                p.font.size = Pt(18)
                p.level = 0

        first_slide = True
        for point in content:
            word_count = len(point.split())
            if current_word_count + word_count <= max_words_per_slide:
                current_content.append(point)
                current_word_count += word_count
            else:
                add_slide(current_content, show_subtitle=first_slide)
                first_slide = False
                current_content = [point]
                current_word_count = word_count

        if current_content:
            add_slide(current_content, show_subtitle=first_slide)

    prs.save(output_file)
    print(f"[✔] PowerPoint saved as: {output_file}")

# اختبار سريع
if __name__ == "__main__":
    title = "Oral Pathology - Sample Lecture"
    slides = [
        {
            "title": "Introduction",
            "subtitle": "Biopsy Principles and Techniques",
            "content": [
                "Oral and maxillofacial pathology is the specialty of dentistry...",
                "Surgical pathology deals with diagnosis by microscopic examination...",
                # أضف نقاط أكثر لتجربة التقسيم
            ]
        },
        {
            "title": "Types of Biopsy",
            "subtitle": "Based on Size and Instruments",
            "content": [
                "Incisional biopsy samples part of the lesion.",
                "Excisional biopsy removes the entire lesion.",
                "Cautery biopsy is least suitable for microscopic interpretation."
            ]
        }
    ]
    create_ppt(title, slides, "sample_lecture.pptx")

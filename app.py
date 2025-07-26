import streamlit as st
import fitz  # PyMuPDF
from PIL import Image
from pptx import Presentation
from pptx.util import Inches
import os

st.set_page_config(page_title="PDF to PowerPoint", page_icon="📊")
st.title("📄➡️📊 PDF إلى PowerPoint")

uploaded_file = st.file_uploader("ارفع ملف PDF", type="pdf")

if uploaded_file:
    # قراءة الملف كـ bytes وفتح المستند
    pdf_bytes = uploaded_file.read()
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")

    st.info(f"📄 عدد الصفحات: {len(doc)}")

    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]  # شريحة فارغة بدون عنوان

    with st.spinner("🚀 جاري تحويل الصفحات إلى شرائح..."):
        for page_num, page in enumerate(doc, start=1):
            pix = page.get_pixmap(dpi=150)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)

            img_path = f"temp_page_{page_num}.png"
            img.save(img_path)

            slide = prs.slides.add_slide(blank_slide_layout)
            slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=prs.slide_width)

            os.remove(img_path)  # حذف الصورة بعد إضافتها للشريحة

    output_name = uploaded_file.name.replace(".pdf", ".pptx")
    prs.save(output_name)

    st.success(f"✅ تم تحويل الملف إلى PowerPoint: {output_name}")

    with open(output_name, "rb") as f:
        st.download_button(
            label="⬇️ تحميل العرض التقديمي",
            data=f,
            file_name=output_name,
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )